import base64
import io
import json
import os
import re
from pathlib import Path
from datetime import datetime

import streamlit as st
import pandas as pd
import fitz
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from groq import Groq
from gtts import gTTS


# ==========================================================
# PAGE CONFIG
# ==========================================================

st.set_page_config(
    page_title="AI Study Buddy",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)


# ==========================================================
# CUSTOM CSS
# ==========================================================

st.markdown("""
<style>

html,
body,
[data-testid="stAppViewContainer"] {
    overflow-x: hidden !important;
    overscroll-behavior-y: none !important;
}

.block-container {
    max-width: 1400px;
    padding-top: 1rem;
    padding-bottom: 4rem;
}

.stButton > button {
    border-radius: 10px;
    min-height: 42px;
    transition: 0.2s;
}

.stButton > button:hover {
    transform: scale(1.01);
}

iframe {
    max-width: 100% !important;
    border: none !important;
}

[data-testid="stSidebar"] {
    overflow-x: hidden !important;
}

[class*="viewerBadge"],
[class*="styles_viewerBadge"],
[data-testid="stAppDeployButton"] {
    display: none !important;
}

</style>
""", unsafe_allow_html=True)


# ==========================================================
# CONSTANTS
# ==========================================================

UPLOAD_EXTENSIONS = [
    "pdf",
    "docx",
    "txt",
    "md",
    "csv",
    "json",
    "xml",
    "html",
    "htm",
    "rtf",
    "xlsx",
    "xlsm",
    "xls",
    "pptx",
    "jpg",
    "jpeg",
    "png",
    "webp",
    "bmp",
    "gif",
    "tif",
    "tiff"
]

IMAGE_EXTENSIONS = {
    "jpg",
    "jpeg",
    "png",
    "webp",
    "bmp",
    "gif",
    "tif",
    "tiff"
}

LANGUAGES = {
    "English": "en",
    "Hindi": "hi",
    "Marathi": "mr",
    "Gujarati": "gu",
    "Bengali": "bn",
    "Tamil": "ta",
    "Telugu": "te",
    "Kannada": "kn",
    "Malayalam": "ml",
    "Punjabi": "pa",
    "Urdu": "ur",
    "French": "fr",
    "Spanish": "es",
    "German": "de",
    "Arabic": "ar"
}


# ==========================================================
# SESSION STATE
# ==========================================================

def init_state():

    defaults = {

        "messages": [],

        "documents": [],

        "file_signature": "",

        "index": None,

        "topics": [],

        "mindmap": None,

        "flashcards": None,

        "quiz": None,

        "quiz_result": None,

        "image_result": None,

        "image_source_name": None,

        "translation": None,

        "generated_image": None,

        "question_count": 0,

        "quiz_count": 0,

        "image_pending": None,

        "last_sources": []

    }

    for key, value in defaults.items():

        if key not in st.session_state:

            st.session_state[key] = value

    Path("data/uploads").mkdir(
        parents=True,
        exist_ok=True
    )

    Path("data/memory").mkdir(
        parents=True,
        exist_ok=True
    )

    Path("data/reports").mkdir(
        parents=True,
        exist_ok=True
    )


init_state()


# ==========================================================
# GROQ API
# ==========================================================

def get_groq_key():

    key = ""

    try:
        key = st.secrets.get(
            "GROQ_API_KEY",
            ""
        )
    except Exception:
        pass

    if not key:

        key = os.getenv(
            "GROQ_API_KEY",
            ""
        )

    if not key:

        raise RuntimeError(
            "GROQ_API_KEY missing. "
            "Please add it to .streamlit/secrets.toml"
        )

    return key


def groq_client():

    return Groq(
        api_key=get_groq_key()
    )


# ==========================================================
# TEXT AI
# ==========================================================

def text_ai(
    prompt,
    temperature=0.2
):

    response = groq_client().chat.completions.create(

        model="openai/gpt-oss-20b",

        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],

        temperature=temperature
    )

    return (
        response.choices[0]
        .message.content
        or ""
    ).strip()


# ==========================================================
# WEB / EXTERNAL INFORMATION
# ==========================================================

def web_ai(prompt):

    response = groq_client().chat.completions.create(

        model="groq/compound",

        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],

        temperature=0.2,

        max_completion_tokens=4096,

        compound_custom={
            "tools": {
                "enabled_tools": [
                    "web_search",
                    "visit_website"
                ]
            }
        }
    )

    message = response.choices[0].message

    answer = (
        message.content
        or ""
    ).strip()

    sources = []

    try:

        executed_tools = (
            getattr(
                message,
                "executed_tools",
                None
            )
            or []
        )

        for tool in executed_tools:

            search_results = (
                getattr(
                    tool,
                    "search_results",
                    None
                )
                or []
            )

            for item in search_results:

                if (
                    isinstance(item, dict)
                    and item.get("url")
                ):

                    sources.append({

                        "title":
                        item.get(
                            "title",
                            "Source"
                        ),

                        "url":
                        item["url"]

                    })

    except Exception:
        pass

    return answer, sources


# ==========================================================
# IMAGE UNDERSTANDING + OCR
# ==========================================================

def vision_ai(
    image_bytes,
    mime_type,
    prompt
):

    encoded = base64.b64encode(
        image_bytes
    ).decode("utf-8")

    response = groq_client().chat.completions.create(

        model="qwen/qwen3.5-27b",

        messages=[
            {
                "role": "user",

                "content": [

                    {
                        "type": "text",
                        "text": prompt
                    },

                    {
                        "type": "image_url",

                        "image_url": {
                            "url":
                            f"data:{mime_type};base64,{encoded}"
                        }
                    }

                ]
            }
        ],

        temperature=0.1
    )

    return (
        response.choices[0]
        .message.content
        or ""
    ).strip()


# ==========================================================
# FILE EXTENSION
# ==========================================================

def get_extension(
    filename
):

    return (
        Path(filename)
        .suffix
        .lower()
        .lstrip(".")
    )


# ==========================================================
# PDF EXTRACTION
# ==========================================================

def extract_pdf(data):

    pdf = fitz.open(
        stream=data,
        filetype="pdf"
    )

    pages = []

    for page_number, page in enumerate(
        pdf,
        start=1
    ):

        text = page.get_text(
            "text"
        ).strip()

        if text:

            pages.append(
                f"[Page {page_number}]\n{text}"
            )

    return "\n\n".join(
        pages
    )


# ==========================================================
# DOCX
# ==========================================================

def extract_docx(data):

    document = Document(
        io.BytesIO(data)
    )

    result = []

    for paragraph in document.paragraphs:

        text = paragraph.text.strip()

        if text:

            result.append(text)

    for table in document.tables:

        for row in table.rows:

            row_text = " | ".join(
                cell.text.strip()
                for cell in row.cells
            )

            if row_text:

                result.append(
                    row_text
                )

    return "\n".join(
        result
    )


# ==========================================================
# POWERPOINT
# ==========================================================

def extract_pptx(data):

    presentation = Presentation(
        io.BytesIO(data)
    )

    slides_text = []

    for number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        texts = []

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                text = shape.text.strip()

                if text:

                    texts.append(
                        text
                    )

        if texts:

            slides_text.append(
                f"[Slide {number}]\n"
                + "\n".join(texts)
            )

    return "\n\n".join(
        slides_text
    )


# ==========================================================
# EXCEL
# ==========================================================

def extract_excel(data):

    excel = pd.ExcelFile(
        io.BytesIO(data)
    )

    sheets = []

    for sheet in excel.sheet_names:

        dataframe = pd.read_excel(
            excel,
            sheet_name=sheet,
            header=None
        )

        dataframe = dataframe.fillna("")

        sheets.append(

            f"[Sheet {sheet}]\n"
            +
            dataframe.to_csv(
                index=False,
                header=False
            )
        )

    return "\n\n".join(
        sheets
    )


# ==========================================================
# TEXT FILES
# ==========================================================

def extract_text(data):

    return data.decode(
        "utf-8",
        errors="ignore"
    )


# ==========================================================
# UNIVERSAL FILE EXTRACTION
# ==========================================================

def extract_file(uploaded_file):

    extension = get_extension(
        uploaded_file.name
    )

    data = uploaded_file.getvalue()

    if extension == "pdf":

        return (
            extract_pdf(data),
            "document"
        )

    if extension == "docx":

        return (
            extract_docx(data),
            "document"
        )

    if extension == "pptx":

        return (
            extract_pptx(data),
            "document"
        )

    if extension in {
        "xlsx",
        "xlsm",
        "xls"
    }:

        return (
            extract_excel(data),
            "document"
        )

    if extension in {
        "txt",
        "md",
        "csv",
        "json",
        "xml",
        "html",
        "htm",
        "rtf"
    }:

        return (
            extract_text(data),
            "document"
        )

    if extension in IMAGE_EXTENSIONS:

        return (
            "",
            "image"
        )

    if extension == "doc":

        return (
            "",
            "legacy_doc"
        )

    return (
        "",
        "unsupported"
    )


# ==========================================================
# TEXT CHUNKS
# ==========================================================

def make_chunks(
    text,
    size=1500,
    overlap=250
):

    text = re.sub(
        r"\s+",
        " ",
        text
    ).strip()

    if not text:

        return []

    chunks = []

    start = 0

    while start < len(text):

        end = min(
            len(text),
            start + size
        )

        if end < len(text):

            cut = text.rfind(
                " ",
                start,
                end
            )

            if cut > start + 500:

                end = cut

        chunks.append(
            text[start:end]
        )

        if end >= len(text):

            break

        start = max(
            0,
            end - overlap
        )

    return chunks


# ==========================================================
# SEARCH INDEX
# ==========================================================

def build_index(documents):

    texts = [
        item["text"]
        for item in documents
        if item.get("text")
    ]

    if not texts:

        return None

    vectorizer = TfidfVectorizer(
        max_features=20000,
        ngram_range=(1, 2),
        stop_words="english"
    )

    matrix = vectorizer.fit_transform(
        texts
    )

    return {

        "vectorizer":
        vectorizer,

        "matrix":
        matrix,

        "documents":
        documents

    }


def search_index(
    index,
    query,
    k=8
):

    if not index:

        return []

    query_vector = (
        index["vectorizer"]
        .transform([query])
    )

    scores = cosine_similarity(
        query_vector,
        index["matrix"]
    )[0]

    indexes = (
        scores
        .argsort()[::-1][:k]
    )

    results = []

    for index_number in indexes:

        score = float(
            scores[index_number]
        )

        if score > 0:

            results.append({

                **index[
                    "documents"
                ][
                    int(index_number)
                ],

                "score":
                score

            })

    return results


def material_context(
    query,
    k=8
):

    results = search_index(
        st.session_state.index,
        query,
        k
    )

    context = "\n\n".join(
        item["text"]
        for item in results
    )

    return (
        context,
        results
    )


# ==========================================================
# MEMORY
# ==========================================================

MEMORY_FILE = Path(
    "data/memory/memories.json"
)


def load_memory():

    try:

        return json.loads(
            MEMORY_FILE.read_text(
                encoding="utf-8"
            )
        )

    except Exception:

        return {}


def save_memory(memory):

    MEMORY_FILE.write_text(
        json.dumps(
            memory,
            ensure_ascii=False,
            indent=2
        ),
        encoding="utf-8"
    )


def capture_memory(text):

    memory = load_memory()

    patterns = [

        r"\bremember(?: that)?\s+"
        r"(?:my\s+)?"
        r"(.+?)\s+(?:is|are)\s+(.+)$",

        r"\bmy\s+"
        r"(.{2,50}?)\s+"
        r"(?:is|are)\s+(.+)$",

        r"\bmera\s+"
        r"(.{2,50}?)\s+"
        r"(?:hai|h|he)\s+(.+)$",

        r"\bmere\s+"
        r"(.{2,50}?)\s+"
        r"(?:ka naam|hai|h)\s+(.+)$"
    ]

    for pattern in patterns:

        match = re.search(
            pattern,
            text.strip(),
            re.IGNORECASE
        )

        if match:

            key = re.sub(
                r"[^\w ]",
                "",
                match.group(1)
                .strip()
                .lower()
            )

            key = key.replace(
                " ",
                "_"
            )

            value = (
                match.group(2)
                .strip()
                .rstrip(".")
            )

            if (
                key
                and len(value) <= 300
            ):

                memory[key] = value

                save_memory(
                    memory
                )

            break

    return memory


def memory_context():

    memory = load_memory()

    if not memory:

        return "No saved memory."

    return "\n".join(

        f"- {key.replace('_', ' ')}: {value}"

        for key, value
        in memory.items()

    )


# ==========================================================
# JSON CLEANER
# ==========================================================

def parse_json(text):

    text = re.sub(
        r"```(?:json)?",
        "",
        text,
        flags=re.IGNORECASE
    )

    text = text.replace(
        "```",
        ""
    ).strip()

    start = text.find("[")
    end = text.rfind("]")

    if (
        start >= 0
        and end > start
    ):

        return json.loads(
            text[
                start:end + 1
            ]
        )

    start = text.find("{")
    end = text.rfind("}")

    if (
        start >= 0
        and end > start
    ):

        return json.loads(
            text[
                start:end + 1
            ]
        )

    raise ValueError(
        "AI returned invalid JSON"
    )


# ==========================================================
# TEXT TO SPEECH
# ==========================================================

def make_audio(
    text,
    language
):

    path = (
        "data/reports/"
        "response.mp3"
    )

    gTTS(
        text=text,
        lang=language
    ).save(path)

    return path


# ==========================================================
# ADD CHAT MESSAGE
# ==========================================================

def add_message(
    role,
    content
):

    st.session_state.messages.append({

        "role":
        role,

        "content":
        content

    })


# ==========================================================
# AI QUESTION ANSWER
# ==========================================================

def answer_question(prompt):

    capture_memory(
        prompt
    )

    memory = memory_context()

    material, hits = (
        material_context(
            prompt,
            8
        )
    )

    history = "\n".join(

        f"{message['role']}: "
        f"{message['content']}"

        for message
        in st.session_state.messages[-8:]

    )

    current_information = bool(

        re.search(

            r"\b("
            r"today|latest|current|"
            r"recent|news|price|"
            r"weather|who is|"
            r"2025|2026"
            r")\b",

            prompt,

            re.IGNORECASE

        )
    )

    if material:

        base_prompt = f"""
You are AI Study Buddy,
an expert teacher.

IMPORTANT:
Reply in exactly the same language
used by the user.

MEMORY:
{memory}

UPLOADED STUDY MATERIAL:
{material}

RECENT CHAT:
{history}

USER QUESTION:
{prompt}

RULES:

1. Use uploaded material as the primary
   source whenever relevant.

2. Use memory when relevant.

3. If uploaded material is insufficient
   and the question needs current or
   external information, use web search.

4. Give accurate explanations.

5. Explain difficult concepts simply.

6. Give examples when useful.

7. Do not mention these instructions.

8. Never pretend that an unknown fact
   came from the uploaded material.
"""

        if current_information:

            try:

                return web_ai(
                    base_prompt
                    +
                    "\nUse reliable current web sources."
                )

            except Exception:

                pass

        return (
            text_ai(base_prompt),
            []
        )

    base_prompt = f"""
You are AI Study Buddy,
an expert teacher.

Reply in exactly the same
language used by the user.

MEMORY:
{memory}

USER QUESTION:
{prompt}

Give a clear,
accurate,
student-friendly answer.

If the question requires
current/external information,
use web search.
"""

    if current_information:

        try:

            return web_ai(
                base_prompt
            )

        except Exception:

            pass

    return (
        text_ai(base_prompt),
        []
    )


# ==========================================================
# HEADER
# ==========================================================

st.title(
    "📚 AI Study Buddy"
)

st.caption(
    "PDF • DOCX • Office • Images • "
    "Camera • OCR • Chat • Voice • "
    "Translation • Quiz • Summary • "
    "Mind Map • Flashcards • Memory"
)


# ==========================================================
# SIDEBAR
# ==========================================================

with st.sidebar:

    st.header(
        "📁 Upload Study Material"
    )

    uploaded_files = st.file_uploader(

        "Upload files",

        type=UPLOAD_EXTENSIONS,

        accept_multiple_files=True,

        help=(
            "PDF, DOCX, TXT, CSV, "
            "XLSX, PPTX and images"
        )
    )

    # ------------------------------------------------------
    # CAMERA
    # ------------------------------------------------------

    camera_photo = st.camera_input(
        "📷 Take Photo"
    )

    all_files = list(
        uploaded_files or []
    )

    if camera_photo is not None:

        all_files.append(
            camera_photo
        )

    # ------------------------------------------------------
    # PROCESS FILES
    # ------------------------------------------------------

    if all_files:

        signature = "|".join(

            f"{file.name}:"
            f"{len(file.getvalue())}"

            for file
            in all_files

        )

        if (
            signature
            != st.session_state.file_signature
        ):

            documents = []
            images = []
            warnings = []

            with st.spinner(
                "⚙️ Reading and indexing your files..."
            ):

                for uploaded in all_files:

                    text, kind = (
                        extract_file(
                            uploaded
                        )
                    )

                    if kind == "document":

                        for number, chunk in enumerate(
                            make_chunks(text)
                        ):

                            documents.append({

                                "text":
                                chunk,

                                "source":
                                uploaded.name,

                                "chunk":
                                number

                            })

                    elif kind == "image":

                        images.append(
                            uploaded
                        )

                    elif kind == "legacy_doc":

                        warnings.append(

                            f"{uploaded.name}: "
                            ".doc is the old Word format. "
                            "Save it as .docx for reliable reading."

                        )

                    elif kind == "unsupported":

                        warnings.append(

                            f"{uploaded.name}: "
                            "this file type is not supported "
                            "for text extraction yet."

                        )

            st.session_state.documents = (
                documents
            )

            st.session_state.index = (
                build_index(
                    documents
                )
            )

            st.session_state.file_signature = (
                signature
            )

            if images:

                st.session_state.image_pending = (
                    images[0]
                )

            if documents:

                st.success(
                    f"✅ {len(all_files)} file(s) "
                    "loaded and searchable"
                )

            for warning in warnings:

                st.warning(
                    warning
                )

    # ======================================================
    # IMAGE ANALYSIS
    # ======================================================

    if (
        st.session_state.get(
            "image_pending"
        ) is not None
    ):

        image = (
            st.session_state.image_pending
        )

        st.image(
            image,
            caption=image.name,
            use_container_width=True
        )

        if st.button(
            "👁️ Analyze Image + OCR",
            use_container_width=True
        ):

            with st.spinner(
                "AI is understanding the image..."
            ):

                try:

                    result = vision_ai(

                        image.getvalue(),

                        image.type,

                        """
Analyze this image carefully.

Use OCR to read visible text.

Also understand the visual content.

Return EXACTLY 5 numbered points.

The 5 points should contain the
most important information visible
in the image.

Do not invent details.

Keep every point useful and concise.

Reply in the same language as the
user's request whenever possible.
"""
                    )

                    st.session_state.image_result = (
                        result
                    )

                    st.session_state.image_source_name = (
                        image.name
                    )

                except Exception as error:

                    st.error(
                        f"Image analysis failed: {error}"
                    )

    # ======================================================
    # SMART STUDY
    # ======================================================

    st.divider()

    st.subheader(
        "🧠 Smart Study Tools"
    )

    # ======================================================
    # QUIZ
    # ======================================================

    question_number = st.slider(

        "Quiz questions",

        min_value=1,

        max_value=20,

        value=5

    )

    if st.button(
        "📝 Generate Advanced Quiz",
        use_container_width=True
    ):

        if not st.session_state.index:

            st.warning(
                "Please upload study material first."
            )

        else:

            with st.spinner(
                "Creating advanced quiz..."
            ):

                try:

                    context, _ = (
                        material_context(
                            "important concepts "
                            "definitions applications "
                            "exam topics",
                            15
                        )
                    )

                    quiz_prompt = f"""

Create exactly
{question_number}
advanced but fair MCQs.

Use ONLY the study material
below.

Questions must be medium-sized,
concept-focused and application-focused.

Each question must have exactly
4 options.

Return ONLY valid JSON.

Format:

[
  {{
    "question": "...",
    "options": [
      "...",
      "...",
      "...",
      "..."
    ],
    "answer": "exact option text",
    "explanation": "short explanation"
  }}
]

The answer must exactly match
one of the options.

STUDY MATERIAL:

{context}
"""

                    raw = text_ai(
                        quiz_prompt,
                        0.1
                    )

                    quiz = parse_json(
                        raw
                    )

                    valid = []

                    for question in quiz:

                        if not isinstance(
                            question,
                            dict
                        ):

                            continue

                        options = question.get(
                            "options",
                            []
                        )

                        answer = question.get(
                            "answer"
                        )

                        if (
                            len(options) == 4
                            and answer in options
                        ):

                            valid.append(
                                question
                            )

                    if len(valid) < question_number:

                        raise ValueError(
                            f"AI returned only "
                            f"{len(valid)} valid questions."
                        )

                    st.session_state.quiz = (
                        valid[:question_number]
                    )

                    st.session_state.quiz_result = None

                    st.success(
                        "✅ Advanced quiz generated!"
                    )

                except Exception as error:

                    st.error(
                        f"Quiz generation failed: {error}"
                    )

    # ======================================================
    # SUMMARY
    # ======================================================

    if st.button(
        "📚 Generate Summary",
        use_container_width=True
    ):

        if not st.session_state.index:

            st.warning(
                "Please upload material first."
            )

        else:

            with st.spinner(
                "AI is preparing your summary..."
            ):

                try:

                    context, _ = (
                        material_context(
                            "complete chapters "
                            "topics definitions "
                            "key points",
                            15
                        )
                    )

                    summary = text_ai(

                        f"""
Create a structured,
exam-focused study summary.

Reply in the user's language.

Include:

- Clear headings
- Important definitions
- Key concepts
- Examples
- Exam-important points
- Easy explanations

Study material:

{context}
"""
                    )

                    add_message(

                        "assistant",

                        "## 📚 Quick Revision Summary\n\n"
                        + summary

                    )

                except Exception as error:

                    st.error(
                        f"Summary failed: {error}"
                    )

    # ======================================================
    # FIND TOPICS
    # ======================================================

    if st.button(
        "📚 Find Study Topics",
        use_container_width=True
    ):

        if not st.session_state.index:

            st.warning(
                "Please upload material first."
            )

        else:

            try:

                context, _ = (
                    material_context(
                        "main topics chapters "
                        "headings concepts",
                        15
                    )
                )

                raw = text_ai(

                    f"""
Find 5 to 10 important study
topics from the material.

Return ONLY a numbered list.

Each topic should be short,
unique and meaningful.

Material:

{context}
"""
                )

                topics = []

                for line in raw.splitlines():

                    topic = re.sub(

                        r"^\s*\d+[.)\-:]\s*",

                        "",

                        line

                    ).strip()

                    topic = re.sub(

                        r"^[-*•]\s*",

                        "",

                        topic

                    ).strip()

                    if (
                        topic
                        and len(topic) <= 100
                        and topic not in topics
                    ):

                        topics.append(
                            topic
                        )

                st.session_state.topics = (
                    topics[:10]
                )

                st.session_state.mindmap = None

                st.success(
                    f"✅ Found {len(topics[:10])} topics"
                )

            except Exception as error:

                st.error(
                    f"Topic detection failed: {error}"
                )

    # ======================================================
    # MIND MAP
    # ======================================================

    if st.session_state.topics:

        selected_topic = st.selectbox(

            "🧠 Select topic for Mind Map",

            st.session_state.topics

        )

        if st.button(
            "🗺️ Generate Mind Map",
            use_container_width=True
        ):

            with st.spinner(
                "Creating topic mind map..."
            ):

                try:

                    context, _ = (
                        material_context(
                            selected_topic,
                            10
                        )
                    )

                    dot = text_ai(

                        f"""
Create ONLY valid Graphviz DOT code.

Topic:

{selected_topic}

Requirements:

- Use digraph G
- Use rankdir=LR
- Maximum 15 nodes
- Short readable labels
- Every label must use double quotes
- Connect the central topic to major concepts
- No markdown
- No explanation

Use this study material:

{context}
""",

                        0.1
                    )

                    dot = re.sub(

                        r"```(?:dot|graphviz)?",

                        "",

                        dot,

                        flags=re.IGNORECASE

                    )

                    dot = dot.replace(
                        "```",
                        ""
                    ).strip()

                    start = dot.lower().find(
                        "digraph"
                    )

                    end = dot.rfind(
                        "}"
                    )

                    if (
                        start < 0
                        or end <= start
                    ):

                        raise ValueError(
                            "AI did not return valid Graphviz."
                        )

                    st.session_state.mindmap = (
                        dot[
                            start:end + 1
                        ]
                    )

                except Exception as error:

                    st.error(
                        f"Mind Map failed: {error}"
                    )

    # ======================================================
    # FLASHCARDS
    # ======================================================

    if st.button(
        "🗂️ Generate 6 Flashcards",
        use_container_width=True
    ):

        if not st.session_state.index:

            st.warning(
                "Please upload material first."
            )

        else:

            with st.spinner(
                "Creating 6 flashcards..."
            ):

                try:

                    context, _ = (
                        material_context(
                            "important terms "
                            "definitions "
                            "key concepts",
                            10
                        )
                    )

                    raw = text_ai(

                        f"""
Create exactly 6 important
study flashcards.

Return ONLY JSON.

Format:

[
  {{
    "term": "...",
    "definition": "..."
  }}
]

Study material:

{context}
""",

                        0.1
                    )

                    cards = parse_json(
                        raw
                    )

                    if len(cards) < 6:

                        raise ValueError(
                            "AI returned fewer than 6 flashcards."
                        )

                    st.session_state.flashcards = (
                        cards[:6]
                    )

                except Exception as error:

                    st.error(
                        f"Flashcard generation failed: {error}"
                    )

    # ======================================================
    # TRANSLATION
    # ======================================================

    st.divider()

    st.subheader(
        "🌐 Translation"
    )

    translation_text = st.text_area(
        "Text to translate"
    )

    target_language = st.selectbox(

        "Target language",

        list(
            LANGUAGES.keys()
        ),

        index=1

    )

    if st.button(
        "Translate Text",
        use_container_width=True
    ):

        if translation_text.strip():

            with st.spinner(
                "Translating..."
            ):

                try:

                    st.session_state.translation = (
                        text_ai(

                            f"""
Translate the following text
accurately into
{target_language}.

Preserve the meaning and
formatting.

Return ONLY the translation.

TEXT:

{translation_text}
"""
                        )
                    )

                except Exception as error:

                    st.error(
                        f"Translation failed: {error}"
                    )

    # ======================================================
    # IMAGE GENERATION
    # ======================================================

    st.divider()

    st.subheader(
        "🎨 Image Generation"
    )

    image_prompt = st.text_area(
        "Describe the image you want"
    )

    if st.button(
        "Generate Image",
        use_container_width=True
    ):

        hf_token = ""

        try:

            hf_token = st.secrets.get(
                "HF_TOKEN",
                ""
            )

        except Exception:

            pass

        if not hf_token:

            hf_token = os.getenv(
                "HF_TOKEN",
                ""
            )

        if not hf_token:

            st.info(
                "Image generation is optional. "
                "Add HF_TOKEN to secrets to enable it."
            )

        elif not image_prompt.strip():

            st.warning(
                "Describe the image first."
            )

        else:

            with st.spinner(
                "Generating image..."
            ):

                try:

                    from huggingface_hub import (
                        InferenceClient
                    )

                    client = InferenceClient(
                        provider="hf-inference",
                        token=hf_token
                    )

                    generated = (
                        client.text_to_image(

                            image_prompt,

                            model=(
                                "black-forest-labs/"
                                "FLUX.1-schnell"
                            )

                        )
                    )

                    st.session_state.generated_image = (
                        generated
                    )

                except Exception as error:

                    st.error(
                        f"Image generation failed: {error}"
                    )

    # ======================================================
    # TEXT TO SPEECH
    # ======================================================

    st.divider()

    st.subheader(
        "🔊 Read Aloud"
    )

    speech_language = st.selectbox(

        "Speech language",

        list(
            LANGUAGES.keys()
        )

    )

    speech_text = st.text_area(
        "Text to speak"
    )

    if st.button(
        "🔊 Play Text",
        use_container_width=True
    ):

        if speech_text.strip():

            try:

                audio_file = make_audio(

                    speech_text,

                    LANGUAGES[
                        speech_language
                    ]

                )

                st.audio(
                    audio_file,
                    format="audio/mp3"
                )

            except Exception as error:

                st.error(
                    f"Text-to-speech failed: {error}"
                )

    # ======================================================
    # CLEAR CHAT
    # ======================================================

    st.divider()

    if st.button(
        "🗑️ Clear Chat",
        use_container_width=True
    ):

        st.session_state.messages = []

        st.rerun()

    # ======================================================
    # MEMORY
    # ======================================================

    if st.button(
        "🧠 Show Saved Memory",
        use_container_width=True
    ):

        st.json(
            load_memory()
        )

    # ======================================================
    # STUDY REPORT
    # ======================================================

    report = (

        "📚 AI STUDY BUDDY "
        "- COMPLETE STUDY REPORT\n\n"

        +

        f"Generated: "
        f"{datetime.now().strftime('%d-%B-%Y %H:%M')}\n\n"

        +

        "=" * 60

        +

        "\n\n"

    )

    for message in (
        st.session_state.messages
    ):

        report += (

            f"{message['role'].upper()}:\n"

            f"{message['content']}\n\n"

            +

            "-" * 60

            +

            "\n\n"

        )

    st.download_button(

        "📥 Download Study Report",

        report,

        "AI_Study_Report.txt",

        "text/plain",

        use_container_width=True

    )


# ==========================================================
# MAIN OUTPUT — IMAGE
# ==========================================================

if st.session_state.image_result:

    st.markdown(
        "## 👁️ Image Analysis — 5 Important Points"
    )

    st.markdown(
        st.session_state.image_result
    )

    try:

        audio_file = make_audio(

            st.session_state.image_result,

            LANGUAGES.get(
                speech_language,
                "en"
            )

        )

        st.audio(
            audio_file,
            format="audio/mp3"
        )

    except Exception:

        pass


# ==========================================================
# MAIN OUTPUT — TRANSLATION
# ==========================================================

if st.session_state.translation:

    st.markdown(
        "## 🌐 Translation"
    )

    st.markdown(
        st.session_state.translation
    )


# ==========================================================
# MAIN OUTPUT — GENERATED IMAGE
# ==========================================================

if st.session_state.generated_image:

    st.markdown(
        "## 🎨 Generated Image"
    )

    st.image(
        st.session_state.generated_image,
        use_container_width=True
    )


# ==========================================================
# MAIN OUTPUT — MIND MAP
# ==========================================================

if st.session_state.mindmap:

    st.markdown(
        "## 🗺️ Topic Mind Map"
    )

    st.graphviz_chart(

        st.session_state.mindmap,

        use_container_width=True

    )


# ==========================================================
# MAIN OUTPUT — FLASHCARDS
# ==========================================================

if st.session_state.flashcards:

    st.markdown(
        "## 🗂️ Interactive Flashcards"
    )

    columns = st.columns(3)

    for i, card in enumerate(
        st.session_state.flashcards
    ):

        with columns[i % 3]:

            st.markdown(
                f"### {i + 1}. "
                f"{card.get('term', '')}"
            )

            with st.expander(
                "Show definition"
            ):

                st.write(
                    card.get(
                        "definition",
                        ""
                    )
                )


# ==========================================================
# MAIN OUTPUT — QUIZ
# ==========================================================

if st.session_state.quiz:

    st.markdown(
        "## 📝 Advanced Quiz"
    )

    with st.form(
        "advanced_quiz_form"
    ):

        answers = []

        for i, question in enumerate(
            st.session_state.quiz
        ):

            st.markdown(
                f"**Q{i + 1}. "
                f"{question['question']}**"
            )

            answer = st.radio(

                "Select one",

                question["options"],

                index=None,

                key=f"quiz_answer_{i}",

                label_visibility="collapsed"

            )

            answers.append(
                answer
            )

        submitted = st.form_submit_button(
            "Submit Answers"
        )

    if submitted:

        correct = 0

        wrong = 0

        not_attempted = 0

        details = []

        for question, answer in zip(

            st.session_state.quiz,

            answers

        ):

            if answer is None:

                not_attempted += 1

                details.append(

                    (
                        "Not Attempted",

                        question["answer"],

                        question.get(
                            "explanation",
                            ""
                        )

                    )

                )

            elif answer == question["answer"]:

                correct += 1

                details.append(

                    (
                        "Correct",

                        question["answer"],

                        question.get(
                            "explanation",
                            ""
                        )

                    )

                )

            else:

                wrong += 1

                details.append(

                    (
                        "Wrong",

                        question["answer"],

                        question.get(
                            "explanation",
                            ""
                        )

                    )

                )

        st.session_state.quiz_result = {

            "correct":
            correct,

            "wrong":
            wrong,

            "not_attempted":
            not_attempted,

            "details":
            details

        }

        st.session_state.quiz_count += 1

    result = (
        st.session_state.quiz_result
    )

    if result:

        total = len(
            st.session_state.quiz
        )

        percentage = round(
            result["correct"]
            * 100
            / total
        )

        st.success(
            f"✅ Correct: "
            f"{result['correct']}"
        )

        st.error(
            f"❌ Wrong: "
            f"{result['wrong']}"
        )

        st.warning(
            f"⚪ Not Attempted: "
            f"{result['not_attempted']}"
        )

        st.info(

            f"📊 Final Score: "
            f"{result['correct']}/{total} "
            f"— {percentage}%"

        )

        for number, detail in enumerate(

            result["details"],

            start=1

        ):

            status, correct_answer, explanation = detail

            st.write(

                f"**Q{number}: "
                f"{status}** — "
                f"Correct answer: "
                f"{correct_answer}. "
                f"{explanation}"

            )


# ==========================================================
# CHAT
# ==========================================================

st.markdown("---")

st.subheader(
    "💬 Chat with your Study Buddy"
)


# ----------------------------------------------------------
# CHAT HISTORY
# ----------------------------------------------------------

for message in (
    st.session_state.messages
):

    with st.chat_message(
        message["role"]
    ):

        st.markdown(
            message["content"]
        )


# ==========================================================
# VOICE INPUT
# ==========================================================

voice_audio = st.audio_input(
    "🎤 Record your question"
)

voice_prompt = None

if voice_audio:

    try:

        transcription = (
            groq_client()
            .audio
            .transcriptions
            .create(

                file=(

                    "voice.wav",

                    voice_audio.getvalue()

                ),

                model=(
                    "whisper-large-v3-turbo"
                ),

                response_format="json"

            )
        )

        voice_prompt = (
            transcription.text
            or ""
        ).strip()

        if voice_prompt:

            st.info(
                f"🎤 You said: "
                f"{voice_prompt}"
            )

    except Exception as error:

        st.error(
            f"Voice-to-text failed: {error}"
        )


# ==========================================================
# CHAT INPUT
# ==========================================================

text_prompt = st.chat_input(
    "Ask anything about your material..."
)

prompt = (
    text_prompt
    or voice_prompt
)


# ==========================================================
# CHAT PROCESSING
# ==========================================================

if prompt:

    add_message(
        "user",
        prompt
    )

    st.session_state.question_count += 1

    with st.chat_message(
        "user"
    ):

        st.markdown(
            prompt
        )

    with st.chat_message(
        "assistant"
    ):

        with st.spinner(
            "Thinking..."
        ):

            try:

                answer, sources = (
                    answer_question(
                        prompt
                    )
                )

                st.markdown(
                    answer
                )

                st.session_state.last_sources = (
                    sources
                )

                if sources:

                    st.markdown(
                        "**🌐 Sources:**"
                    )

                    for source in sources[:6]:

                        st.markdown(

                            f"- "
                            f"[{source['title']}]"
                            f"({source['url']})"

                        )

                # ------------------------------------------------
                # AI ANSWER SPEECH
                # ------------------------------------------------

                try:

                    language = LANGUAGES.get(
                        speech_language,
                        "en"
                    )

                    audio_file = make_audio(

                        answer,

                        language

                    )

                    st.audio(
                        audio_file,
                        format="audio/mp3"
                    )

                except Exception:

                    pass

                add_message(
                    "assistant",
                    answer
                )

            except Exception as error:

                error_message = (
                    f"❌ AI Error: {error}"
                )

                st.error(
                    error_message
                )

                add_message(
                    "assistant",
                    error_message
                )


# ==========================================================
# ANALYTICS
# ==========================================================

with st.sidebar:

    st.divider()

    st.subheader(
        "📊 My Progress"
    )

    st.metric(
        "Questions",
        st.session_state.question_count
    )

    file_count = 0

    if st.session_state.file_signature:

        file_count = len(
            st.session_state.file_signature.split("|")
        )

    st.metric(
        "Files",
        file_count
    )

    words_learned = sum(

        len(
            message["content"].split()
        )

        for message
        in st.session_state.messages

        if message["role"] == "assistant"

    )

    st.metric(
        "Words Learned",
        words_learned
    )

    st.progress(

        min(
            st.session_state.question_count * 10,
            100
        )

    )    
